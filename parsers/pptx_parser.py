"""PPTX parser using python-pptx to extract sections from slides."""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches  # noqa: F401 (imported for availability)

from config.settings import settings
from parsers.base import BaseParser, ParsedSection
from parsers.image_utils import bytes_to_base64_image

logger = logging.getLogger(__name__)


def _get_slide_title(slide, slide_number: int) -> str:
    """Extract the title from a slide using multiple fallback strategies.

    Priority:
      1. The slide's title placeholder text
      2. The first text shape with content
      3. "Slide N"
    """
    # Strategy 1: title placeholder
    if slide.shapes.title is not None:
        title_text = slide.shapes.title.text.strip()
        if title_text:
            return title_text

    # Strategy 2: first non-empty text shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # Use first line only to keep title short
                first_line = text.split("\n")[0].strip()
                if first_line:
                    return first_line

    # Strategy 3: fallback
    return f"Slide {slide_number}"


def _get_slide_body(slide) -> str:
    """Concatenate text from all non-title shapes on a slide."""
    title_shape = slide.shapes.title
    parts: list[str] = []

    for shape in slide.shapes:
        # Skip the title shape (already captured as section title)
        if shape is title_shape:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)

    return "\n\n".join(parts)


def _extract_slide_images(slide) -> list[dict]:
    """Extract images from slide shapes.

    Iterates over shapes looking for embedded images. Each image is resized
    and converted to a base64 dict compatible with ParsedSection.images.

    Returns:
        List of image dicts: {"data": base64_str, "media_type": str, "caption": str}
    """
    if not settings.EXTRACT_IMAGES:
        return []
    images: list[dict] = []
    for shape in slide.shapes:
        if len(images) >= settings.MAX_IMAGES_PER_SECTION:
            break
        try:
            if hasattr(shape, "image") and shape.image is not None:
                blob = shape.image.blob
                content_type = shape.image.content_type  # e.g. "image/png"
                # Derive file extension from content_type
                ext = "." + content_type.split("/")[-1] if "/" in content_type else ".png"
                b64, media_type = bytes_to_base64_image(
                    blob, ext, settings.IMAGE_MAX_EDGE
                )
                caption = shape.name or "Slide image"
                images.append({
                    "data": b64,
                    "media_type": media_type,
                    "caption": caption,
                })
        except Exception as e:
            logger.debug("Skipping shape image extraction: %s", e)
            continue
    return images


class PptxParser(BaseParser):
    """Parse .pptx files, treating each slide as a section."""

    def supported_extensions(self) -> list[str]:
        return [".pptx"]

    def parse(self, file_path: str) -> list[ParsedSection]:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {file_path}")
        if path.suffix.lower() != ".pptx":
            raise ValueError(f"Not a PPTX file: {file_path}")

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise RuntimeError(f"Failed to open PPTX '{file_path}': {exc}") from exc

        sections: list[ParsedSection] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            title = _get_slide_title(slide, slide_idx)
            body = _get_slide_body(slide)

            # Extract images from this slide
            slide_images = _extract_slide_images(slide)

            sections.append(ParsedSection(
                title=title,
                text=body,
                level=1,
                page_number=slide_idx,
                images=slide_images,
            ))

        return sections
